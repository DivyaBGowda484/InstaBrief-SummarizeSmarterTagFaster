from fastapi import APIRouter, HTTPException, Depends, UploadFile, File, Form, Query
from fastapi.responses import FileResponse
from typing import List, Optional
from datetime import datetime
from bson import ObjectId
import os
import aiofiles
import magic
from pathlib import Path

from app.models.schemas import (
    DocumentCreate, DocumentPublic, DocumentUpdate, 
    FileUploadResponse, ProcessingStatus, FileType
)
from app.services.summarizer import SummarizerService
from app.services.nlp import NLPService
from app.services.tts import TTSService
from config.db import get_database, get_elasticsearch
from app.routes.auth import get_current_user

router = APIRouter()
summarizer = SummarizerService()
nlp_service = NLPService()
tts_service = TTSService()

# Ensure upload directory exists
UPLOAD_DIR = Path("uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

def get_file_type(filename: str) -> FileType:
    """Determine file type from filename"""
    ext = filename.lower().split('.')[-1]
    if ext == 'pdf':
        return FileType.PDF
    elif ext == 'docx':
        return FileType.DOCX
    elif ext == 'pptx':
        return FileType.PPTX
    elif ext == 'txt':
        return FileType.TXT
    else:
        return FileType.TXT

async def extract_text_from_file(file_path: str, file_type: FileType) -> str:
    """Extract text content from uploaded file"""
    try:
        if file_type == FileType.PDF:
            import PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text()
                return text
        
        elif file_type == FileType.DOCX:
            from docx import Document
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        
        elif file_type == FileType.PPTX:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        elif file_type == FileType.TXT:
            async with aiofiles.open(file_path, 'r', encoding='utf-8') as file:
                return await file.read()
        
        else:
            return ""
    
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error extracting text: {str(e)}")

@router.post("/upload", response_model=FileUploadResponse)
async def upload_file(
    file: UploadFile = File(...),
    current_user: dict = Depends(get_current_user)
):
    """Upload a file for processing"""
    
    # Validate file type
    file_type = get_file_type(file.filename)
    if file_type not in [FileType.PDF, FileType.DOCX, FileType.PPTX, FileType.TXT]:
        raise HTTPException(status_code=400, detail="Unsupported file type")
    
    # Validate file size
    content = await file.read()
    if len(content) > 25 * 1024 * 1024:  # 25MB limit
        raise HTTPException(status_code=400, detail="File too large")
    
    # Save file
    file_id = str(ObjectId())
    file_path = UPLOAD_DIR / f"{file_id}_{file.filename}"
    
    async with aiofiles.open(file_path, 'wb') as f:
        await f.write(content)
    
    # Extract text content
    try:
        content_text = await extract_text_from_file(str(file_path), file_type)
    except Exception as e:
        # Clean up file on error
        os.remove(file_path)
        raise HTTPException(status_code=400, detail=f"Error processing file: {str(e)}")
    
    return FileUploadResponse(
        file_id=file_id,
        filename=file.filename,
        file_type=file_type.value,
        file_size=len(content),
        processing_status=ProcessingStatus.PENDING
    )

@router.post("/process")
async def process_document(
    title: str = Form(...),
    content: str = Form(...),
    file_type: str = Form(...),
    language: str = Form("en"),
    algorithm: str = Form("textrank"),
    max_length: int = Form(150),
    current_user: dict = Depends(get_current_user)
):
    """Process a document with AI summarization and tagging"""
    
    try:
        # Initialize services
        await summarizer.initialize()
        await nlp_service.initialize()
        
        # Generate summary
        summary_result = await summarizer.generate_summary(
            content, 
            max_length=max_length, 
            algorithm=algorithm
        )
        
        # Extract tags and entities
        keywords = await nlp_service.extract_keywords(content, max_keywords=10)
        entities = await nlp_service.extract_entities(content)
        sentiment = await nlp_service.analyze_sentiment(content)
        language_detected = await nlp_service.detect_language(content)
        
        # Prepare document data
        document_data = {
            "title": title,
            "content": content,
            "summary": summary_result["summary"],
            "file_type": file_type,
            "language": language_detected,
            "sentiment": sentiment["compound"],
            "tags": [kw["word"] for kw in keywords[:5]],
            "entities": [ent["text"] for ent in entities[:10]],
            "processing_status": ProcessingStatus.COMPLETED,
            "algorithm_used": summary_result["algorithm_used"],
            "processing_time": summary_result["processing_time"],
            "compression_ratio": summary_result["compression_ratio"],
            "user_id": current_user["id"],
            "created_at": datetime.utcnow(),
            "updated_at": datetime.utcnow()
        }
        
        # Save to database
        db = await get_database()
        result = await db.documents.insert_one(document_data)
        document_id = str(result.inserted_id)
        
        # Index in Elasticsearch
        try:
            es = await get_elasticsearch()
            await es.index(
                index="instabrief_documents",
                id=document_id,
                body={
                    "title": title,
                    "content": content,
                    "summary": summary_result["summary"],
                    "tags": document_data["tags"],
                    "entities": document_data["entities"],
                    "user_id": current_user["id"],
                    "created_at": datetime.utcnow().isoformat(),
                    "file_type": file_type,
                    "language": language_detected,
                    "sentiment": sentiment["compound"]
                }
            )
        except Exception as e:
            print(f"Elasticsearch indexing failed: {e}")
        
        return {
            "id": document_id,
            "summary": summary_result["summary"],
            "tags": document_data["tags"],
            "entities": document_data["entities"],
            "sentiment": sentiment,
            "processing_time": summary_result["processing_time"],
            "compression_ratio": summary_result["compression_ratio"]
        }
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")

@router.get("/", response_model=List[DocumentPublic])
async def get_documents(
    skip: int = Query(0, ge=0),
    limit: int = Query(20, ge=1, le=100),
    q: Optional[str] = Query(None),
    type: Optional[str] = Query(None),
    sort: str = Query("newest"),
    current_user: dict = Depends(get_current_user)
):
    """Get user's documents with optional search and filtering"""
    
    db = await get_database()
    
    # Build query
    query = {"user_id": current_user["id"]}
    
    if type and type != "all":
        query["file_type"] = type
    
    if q:
        # Text search in title, content, and summary
        query["$or"] = [
            {"title": {"$regex": q, "$options": "i"}},
            {"content": {"$regex": q, "$options": "i"}},
            {"summary": {"$regex": q, "$options": "i"}},
            {"tags": {"$in": [{"$regex": q, "$options": "i"}]}}
        ]
    
    # Build sort
    sort_field = "created_at"
    sort_direction = -1
    if sort == "oldest":
        sort_direction = 1
    elif sort == "name":
        sort_field = "title"
        sort_direction = 1
    elif sort == "size":
        sort_field = "file_size"
        sort_direction = -1
    
    # Execute query
    cursor = db.documents.find(query).sort(sort_field, sort_direction).skip(skip).limit(limit)
    documents = []
    
    async for doc in cursor:
        doc["id"] = str(doc["_id"])
        documents.append(DocumentPublic(**doc))
    
    return documents

@router.get("/{document_id}", response_model=DocumentPublic)
async def get_document(
    document_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Get a specific document by ID"""
    
    try:
        db = await get_database()
        doc = await db.documents.find_one({
            "_id": ObjectId(document_id),
            "user_id": current_user["id"]
        })
        
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        doc["id"] = str(doc["_id"])
        return DocumentPublic(**doc)
    
    except Exception:
        raise HTTPException(status_code=404, detail="Document not found")

@router.put("/{document_id}", response_model=DocumentPublic)
async def update_document(
    document_id: str,
    document_update: DocumentUpdate,
    current_user: dict = Depends(get_current_user)
):
    """Update a document"""
    
    try:
        db = await get_database()
        
        # Check if document exists and belongs to user
        doc = await db.documents.find_one({
            "_id": ObjectId(document_id),
            "user_id": current_user["id"]
        })
        
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Update document
        update_data = document_update.dict(exclude_unset=True)
        update_data["updated_at"] = datetime.utcnow()
        
        await db.documents.update_one(
            {"_id": ObjectId(document_id)},
            {"$set": update_data}
        )
        
        # Get updated document
        updated_doc = await db.documents.find_one({"_id": ObjectId(document_id)})
        updated_doc["id"] = str(updated_doc["_id"])
        
        return DocumentPublic(**updated_doc)
    
    except Exception:
        raise HTTPException(status_code=404, detail="Document not found")

@router.delete("/{document_id}")
async def delete_document(
    document_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Delete a document"""
    
    try:
        db = await get_database()
        
        # Check if document exists and belongs to user
        doc = await db.documents.find_one({
            "_id": ObjectId(document_id),
            "user_id": current_user["id"]
        })
        
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Delete from database
        await db.documents.delete_one({"_id": ObjectId(document_id)})
        
        # Delete from Elasticsearch
        try:
            es = await get_elasticsearch()
            await es.delete(index="instabrief_documents", id=document_id)
        except Exception:
            pass  # Ignore Elasticsearch errors
        
        return {"message": "Document deleted successfully"}
    
    except Exception:
        raise HTTPException(status_code=404, detail="Document not found")

@router.get("/{document_id}/export")
async def export_document(
    document_id: str,
    current_user: dict = Depends(get_current_user)
):
    """Export document summary as text file"""
    
    try:
        db = await get_database()
        doc = await db.documents.find_one({
            "_id": ObjectId(document_id),
            "user_id": current_user["id"]
        })
        
        if not doc:
            raise HTTPException(status_code=404, detail="Document not found")
        
        # Create export content
        content = f"Document: {doc['title']}\n"
        content += f"Created: {doc['created_at'].strftime('%Y-%m-%d %H:%M:%S')}\n"
        content += f"File Type: {doc.get('file_type', 'Unknown')}\n"
        content += f"Language: {doc.get('language', 'Unknown')}\n"
        content += f"Algorithm: {doc.get('algorithm_used', 'Unknown')}\n\n"
        content += f"SUMMARY:\n{doc.get('summary', 'No summary available')}\n\n"
        
        if doc.get('tags'):
            content += f"TAGS: {', '.join(doc['tags'])}\n\n"
        
        if doc.get('entities'):
            content += f"ENTITIES: {', '.join(doc['entities'])}\n\n"
        
        content += f"ORIGINAL CONTENT:\n{doc.get('content', '')}"
        
        # Save to temporary file
        temp_file = f"temp_export_{document_id}.txt"
        async with aiofiles.open(temp_file, 'w', encoding='utf-8') as f:
            await f.write(content)
        
        return FileResponse(
            temp_file,
            media_type='text/plain',
            filename=f"{doc['title']}_summary.txt"
        )
    
    except Exception:
        raise HTTPException(status_code=404, detail="Document not found")
